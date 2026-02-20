
import io
import pandas as pd
import streamlit as st
import numpy as np
from datetime import datetime

# --- PPTX Generation Utilities (New) ---
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_premium_pptx(config, df, kpis, findings, recommendations, charts):
    """
    Generate a premium PowerPoint presentation with a coherent narrative flow.
    """
    prs = Presentation()
    
    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 23, 42) # Dark Navy
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = config.get('title', 'Executive Report')
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if config.get('subtitle'):
        p = title_frame.add_paragraph()
        p.text = str(config.get('subtitle')) # Ensure string
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(148, 163, 184) # Slate 400
        p.alignment = PP_ALIGN.CENTER
        
    # Meta
    meta_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
    meta_frame = meta_box.text_frame
    p = meta_frame.paragraphs[0]
    company = config.get('company', 'Plotiva')
    date_str = datetime.now().strftime('%B %Y')
    p.text = f"{company} | {date_str}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(148, 163, 184)
    p.alignment = PP_ALIGN.CENTER

    # --- Helper: Add Slide Header ---
    def add_slide_header(slide, text):
        # Top Bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.0))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(15, 118, 110) # Teal
        bar.line.fill.background()
        
        # Text
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

    # 2. Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Executive Summary")
    
    # Summary Text
    summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
    tf = summary_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = config.get('summary', 'No summary provided.')
    p.font.size = Pt(18)
    p.line_spacing = 1.2

    # Highlights Box
    if findings:
        hl_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.0), Inches(9), Inches(3))
        hl_box.fill.solid()
        hl_box.fill.fore_color.rgb = RGBColor(240, 253, 250) # Light Teal
        hl_box.line.color.rgb = RGBColor(20, 184, 166)
        
        hl_tf = hl_box.text_frame
        hl_tf.margin_left = Inches(0.2)
        hl_tf.margin_top = Inches(0.2)
        
        p = hl_tf.paragraphs[0]
        p.text = "Key Strategic Insights"
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(15, 118, 110)
        
        for f in findings[:4]: # Limit to 4
            p = hl_tf.add_paragraph()
            # Clean HTML tags if present
            clean_text = f.get('text', '').replace('<b>', '').replace('</b>', '').replace('<br/>', ': ')
            p.text = f"• {clean_text}"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(51, 65, 85)
            p.level = 0
            p.space_before = Pt(6)

    # 3. KPI Dashboard
    if kpis:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_header(slide, "Performance Snapshot")
        
        # Grid layout
        margin_x = 0.5
        margin_y = 1.5
        gap = 0.5
        card_w = 4.0
        card_h = 2.0
        
        for i, kpi in enumerate(kpis[:4]): # Max 4
            row = i // 2
            col = i % 2
            x = margin_x + col * (card_w + gap)
            y = margin_y + row * (card_h + gap)
            
            # Card BG
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(card_h))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(203, 213, 225)
            card.shadow.inherit = True
            
            # Label
            lbl = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.2), Inches(card_w-0.4), Inches(0.5))
            lp = lbl.text_frame.paragraphs[0]
            lp.text = list(kpi.values())[0].upper() if isinstance(kpi, dict) else "METRIC" # Robustness check
            if isinstance(kpi, dict) and 'name' in kpi: lp.text = kpi['name'].upper()
            lp.font.color.rgb = RGBColor(100, 116, 139)
            lp.font.size = Pt(12)
            
            # Value
            val = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.6), Inches(card_w-0.4), Inches(0.8))
            vp = val.text_frame.paragraphs[0]
            vp.text = str(kpi.get('value', '0'))
            vp.font.bold = True
            vp.font.size = Pt(36)
            vp.font.color.rgb = RGBColor(15, 23, 42)
            
            # Change
            chg = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+1.4), Inches(card_w-0.4), Inches(0.4))
            cp = chg.text_frame.paragraphs[0]
            cp.text = str(kpi.get('change', '-'))
            cp.font.size = Pt(14)
            if 'success' in kpi.get('status', ''):
                cp.font.color.rgb = RGBColor(34, 197, 94) # Green
            elif 'danger' in kpi.get('status', ''):
                cp.font.color.rgb = RGBColor(239, 68, 68) # Red
            else:
                cp.font.color.rgb = RGBColor(100, 116, 139)

    # 4. Visual Analysis (One chart per slide)
    for i, chart_info in enumerate(charts):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_header(slide, f"Visual Analysis: {chart_info.get('name', f'Chart {i+1}')}")
        
        # Render chart to image
        try:
            fig = chart_info.get('figure')
            if fig:
                img_bytes = fig.to_image(format='png', width=1000, height=550, scale=2.0)
                slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(0.5), Inches(1.5), width=Inches(9))
                
                # Insight Caption
                cap_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.8), Inches(9), Inches(0.6))
                cap_box.fill.solid()
                cap_box.fill.fore_color.rgb = RGBColor(240, 253, 250)
                cap_box.line.color.rgb = RGBColor(20, 184, 166)
                
                tf = cap_box.text_frame
                tf.margin_left = Inches(0.1)
                p = tf.paragraphs[0]
                ins_text = chart_info.get('caption', '')
                if not ins_text or len(ins_text) < 10: ins_text = "Analysis reveals key distribution patterns."
                p.text = f"💡 INSIGHT: {ins_text}"
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(15, 118, 110)
        except Exception:
            pass

    # 5. Recommendations
    if recommendations:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_header(slide, "Strategic Recommendations")
        
        y_pos = 1.5
        for rec in recommendations[:3]: # Max 3 fits well
            # Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.4))
            tp = title_box.text_frame.paragraphs[0]
            tp.text = f"🚀 {rec.get('title', 'Action Item')}"
            tp.font.bold = True
            tp.font.size = Pt(18)
            tp.font.color.rgb = RGBColor(15, 23, 42)
            
            # Desc
            desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos+0.5), Inches(8.5), Inches(1.0))
            dp = desc_box.text_frame.paragraphs[0]
            dp.text = rec.get('description', '')
            dp.font.size = Pt(14)
            dp.font.color.rgb = RGBColor(71, 85, 105)
            
            y_pos += 1.8

    # 6. Appendix / End
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 23, 42)
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return prs
